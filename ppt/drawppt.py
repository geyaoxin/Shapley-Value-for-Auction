
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
# 创建幻灯片 ------
prs = Presentation()


print(dir(prs))
print(prs.slide_height)
#10 inches * 7.5 inches



graph=prs.slide_layouts[1]
slide=prs.slides.add_slide(graph)
left = top = width = height = Inches(1.0)

slide.shapes.add_shape(
    MSO_SHAPE.OVAL, left, top, width, height
)

for shape in slide.shapes:
    print(shape.shape_type)
prs.save("test.pptx")